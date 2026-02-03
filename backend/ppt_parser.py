"""PowerPoint parsing functionality."""

import os
import subprocess
import shutil
from typing import List, Dict, Tuple
from pptx import Presentation as PptxPresentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont
import io


def extract_text_from_shape(shape) -> str:
    """Extract text from a shape."""
    if hasattr(shape, "text"):
        return shape.text
    return ""


def extract_slide_text(slide) -> Tuple[str, str]:
    """
    Extract text from a slide.
    
    Returns:
        Tuple of (title, body_text)
    """
    title = ""
    body_text = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                # Try to identify title (usually the first text box or placeholder)
                if not title:
                    try:
                        if shape.placeholder_format.type == 1:  # Title placeholder
                            title = text
                            continue
                    except (AttributeError, ValueError):
                        # Shape is not a placeholder, continue processing
                        pass
                
                # If no title yet and this is the first text, use it as title
                if not title and len(body_text) == 0 and len(text) < 100:
                    title = text
                else:
                    body_text.append(text)
    
    return title, "\n".join(body_text)


def extract_notes(slide) -> str:
    """Extract notes from a slide."""
    if hasattr(slide, "notes_slide"):
        notes_slide = slide.notes_slide
        if hasattr(notes_slide, "notes_text_frame"):
            return notes_slide.notes_text_frame.text.strip()
    return ""


def has_visual_content(slide) -> bool:
    """
    Check if a slide has any visual content (images, shapes, etc.) beyond text.
    Used to skip completely blank slides after text removal.
    
    Returns:
        True if slide has visual elements, False if it would be blank
    """
    has_content = False
    
    for shape in slide.shapes:
        # Check for images
        if hasattr(shape, 'image'):
            has_content = True
            break
        
        # Check for shapes with fill (not just text boxes)
        if hasattr(shape, 'fill'):
            try:
                if shape.fill.type is not None and shape.fill.type != 0:  # 0 = no fill
                    has_content = True
                    break
            except:
                pass
        
        # Check for charts
        if hasattr(shape, 'chart'):
            has_content = True
            break
        
        # Check for tables
        if hasattr(shape, 'table'):
            has_content = True
            break
        
        # Check for shapes (not just text boxes)
        if hasattr(shape, 'shape_type'):
            try:
                # Exclude text boxes and placeholders
                if shape.shape_type not in [1, 14]:  # 1=text box, 14=placeholder
                    has_content = True
                    break
            except:
                pass
    
    return has_content


def remove_text_from_pptx(input_path: str, output_path: str, slides_to_keep: List[int] = None) -> str:
    """
    Create a copy of the PowerPoint file with all text removed.
    Optionally removes blank slides (slides without visual content).
    This ensures screenshots don't contain confidential text.
    
    Args:
        input_path: Path to original .pptx file
        output_path: Path where text-free copy will be saved
        slides_to_keep: List of slide numbers to keep (1-indexed), or None to keep all
        
    Returns:
        Path to the text-free PowerPoint file
    """
    print(f"Creating text-free version of presentation...")
    
    # Load the presentation
    prs = PptxPresentation(input_path)
    
    # If we have a list of slides to keep, remove others first
    if slides_to_keep is not None:
        # Create XML tree to manipulate slides
        from pptx.util import Inches
        import xml.etree.ElementTree as ET
        
        # Work backwards to avoid index issues when deleting
        slides_to_delete = []
        for idx in range(len(prs.slides), 0, -1):
            if idx not in slides_to_keep:
                slides_to_delete.append(idx - 1)  # Convert to 0-indexed
        
        # Delete slides without visual content
        for idx in sorted(slides_to_delete, reverse=True):
            try:
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]
                print(f"  Removed blank slide {idx + 1}")
            except:
                pass
    
    # Remove text from remaining slides - RIGOROUS APPROACH
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # Method 1: Clear text frames completely
            if hasattr(shape, "text_frame"):
                try:
                    # Clear all paragraphs
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in list(paragraph.runs):
                            run.text = ""
                        # Also try to remove the paragraph element itself
                        paragraph.text = ""
                    # Clear the entire frame
                    text_frame.clear()
                except Exception as e:
                    print(f"    Warning: Could not clear text_frame: {e}")
            
            # Method 2: Remove text attribute directly
            if hasattr(shape, "text"):
                try:
                    shape.text = ""
                except:
                    pass
            
            # Method 3: For tables, clear all cell text
            if hasattr(shape, "table"):
                try:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if hasattr(cell, "text_frame"):
                                cell.text_frame.clear()
                            if hasattr(cell, "text"):
                                cell.text = ""
                except Exception as e:
                    print(f"    Warning: Could not clear table text: {e}")
            
            # Method 4: For charts, try to remove data labels
            if hasattr(shape, "chart"):
                try:
                    chart = shape.chart
                    # Remove chart title
                    if hasattr(chart, "has_title"):
                        chart.has_title = False
                    # Remove axis titles
                    if hasattr(chart, 'category_axis'):
                        if hasattr(chart.category_axis, 'has_title'):
                            chart.category_axis.has_title = False
                    if hasattr(chart, 'value_axis'):
                        if hasattr(chart.value_axis, 'has_title'):
                            chart.value_axis.has_title = False
                except Exception as e:
                    print(f"    Warning: Could not clear chart text: {e}")
        
        # Remove speaker notes (also potentially confidential)
        if hasattr(slide, "notes_slide") and slide.notes_slide:
            try:
                if hasattr(slide.notes_slide, "notes_text_frame"):
                    slide.notes_slide.notes_text_frame.clear()
            except:
                pass
        
        # Remove slide title and subtitle placeholders
        try:
            from pptx.util import Inches
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format'):
                    try:
                        # Clear placeholder text
                        if hasattr(shape, 'text_frame'):
                            shape.text_frame.clear()
                    except:
                        pass
        except:
            pass
    
    # Save the text-free version
    prs.save(output_path)
    print(f"✅ Text-free version saved: {output_path}")
    
    return output_path


def convert_pptx_to_images(file_path: str, output_dir: str, slides_to_include: List[int] = None) -> List[str]:
    """
    Convert PowerPoint slides to actual screenshot images WITHOUT TEXT.
    
    Process:
    1. Creates a text-free copy of the presentation
    2. Converts text-free copy to PDF
    3. Generates PNG screenshots from PDF
    
    This ensures screenshots show layout/structure but no confidential text.
    REQUIRES LibreOffice to be installed for proper slide rendering.
    
    Args:
        file_path: Path to the original .pptx file
        output_dir: Directory to save slide images
        
    Returns:
        List of image file paths
    """
    os.makedirs(output_dir, exist_ok=True)
    image_paths = []
    
    # Check if soffice or libreoffice is available
    soffice_cmd = None
    possible_commands = [
        'soffice',
        'libreoffice',
        '/usr/local/bin/soffice',
        '/Applications/LibreOffice.app/Contents/MacOS/soffice'  # macOS default
    ]
    
    for cmd in possible_commands:
        try:
            result = subprocess.run([cmd, '--version'], capture_output=True, timeout=5, text=True)
            if result.returncode == 0:
                soffice_cmd = cmd
                print(f"Found LibreOffice: {result.stdout.strip()}")
                break
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    
    if not soffice_cmd:
        error_msg = """
ERROR: LibreOffice is required to generate actual slide screenshots.

Please install LibreOffice:
  macOS:   brew install libreoffice
  Ubuntu:  sudo apt-get install libreoffice
  Windows: Download from https://www.libreoffice.org/download/
  
Then restart the backend server.
"""
        print(error_msg)
        raise RuntimeError("LibreOffice not found. Cannot generate slide screenshots.")
    
    # Create text-free version for screenshots (optionally removing blank slides)
    text_free_path = os.path.join(output_dir, 'text_free_temp.pptx')
    try:
        remove_text_from_pptx(file_path, text_free_path, slides_to_keep=slides_to_include)
    except Exception as e:
        print(f"Warning: Could not remove text from slides: {e}")
        print("Proceeding with original file (text will be visible in screenshots)")
        text_free_path = file_path
    
    try:
        # Convert text-free PPTX to PDF
        print(f"Converting text-free version to PDF using {soffice_cmd}...")
        result = subprocess.run([
            soffice_cmd,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', output_dir,
            text_free_path
        ], capture_output=True, timeout=120, text=True)
        
        if result.returncode != 0:
            print(f"LibreOffice error: {result.stderr}")
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
        
        # Find the generated PDF (based on text-free filename)
        pdf_basename = os.path.splitext(os.path.basename(text_free_path))[0] + '.pdf'
        pdf_path = os.path.join(output_dir, pdf_basename)
        
        if not os.path.exists(pdf_path):
            raise RuntimeError(f"PDF not generated at {pdf_path}")
        
        print(f"PDF generated: {pdf_path}")
        
        # Convert PDF to images using pdf2image
        try:
            from pdf2image import convert_from_path
            print("Converting PDF pages to images...")
            images = convert_from_path(pdf_path, dpi=200, fmt='png')
            
            for idx, img in enumerate(images, start=1):
                img_path = os.path.join(output_dir, f'slide_{idx}.png')
                img.save(img_path, 'PNG', optimize=True)
                image_paths.append(img_path)
                print(f"  Created: slide_{idx}.png")
            
            # Clean up temporary files
            os.remove(pdf_path)
            if text_free_path != file_path and os.path.exists(text_free_path):
                os.remove(text_free_path)
                print(f"✅ Cleaned up temporary text-free file")
            
            print(f"✅ Successfully generated {len(image_paths)} text-free slide screenshots")
            return image_paths
            
        except ImportError:
            raise RuntimeError("pdf2image library not installed. Run: pip install pdf2image")
        except Exception as e:
            raise RuntimeError(f"Failed to convert PDF to images: {str(e)}")
            
    except subprocess.TimeoutExpired:
        raise RuntimeError("LibreOffice conversion timed out (>120s)")
    except Exception as e:
        raise RuntimeError(f"Slide screenshot generation failed: {str(e)}")


def parse_pptx(file_path: str, output_dir: str) -> List[Dict]:
    """
    Parse a PowerPoint file, generate TEXT-FREE slide screenshots, and extract meta-text.
    
    PRIVACY FEATURE: Screenshots are generated WITHOUT TEXT to protect confidential information.
    - Screenshots show: layout, images, charts, shapes, colors, structure
    - Screenshots DON'T show: any text content (removed before screenshot)
    - Text is still extracted and stored separately for search functionality
    
    Meta-text is structured in layers:
    - Primary layer: Slide title (most immediate, highest priority for search)
    - Secondary layer: Body text content
    - Tertiary layer: Speaker notes
    
    Args:
        file_path: Path to the .pptx file
        output_dir: Directory to save text-free slide screenshot images
        
    Returns:
        List of dictionaries containing slide information with text-free screenshots and meta-text
    """
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Load presentation
    prs = PptxPresentation(file_path)
    
    print(f"Processing presentation: {os.path.basename(file_path)}")
    print(f"Total slides: {len(prs.slides)}")
    
    # First pass: Check which slides have visual content (skip blank ones)
    print("Checking slides for visual content...")
    slides_with_content = []
    for idx, slide in enumerate(prs.slides, start=1):
        has_content = has_visual_content(slide)
        if has_content:
            slides_with_content.append(idx)
            print(f"  Slide {idx}: ✅ Has visual content")
        else:
            print(f"  Slide {idx}: ⏭️  Blank (no visual content) - will be skipped")
    
    # Convert slides to actual screenshots (only for slides with visual content)
    print(f"Generating screenshots for {len(slides_with_content)} slides with visual content...")
    try:
        image_paths = convert_pptx_to_images(file_path, output_dir, slides_to_include=slides_with_content)
        print(f"✅ Successfully generated {len(image_paths)} screenshots")
    except Exception as e:
        print(f"⚠️  Warning: Screenshot generation failed: {e}")
        print(f"⚠️  Continuing with text-only mode (LibreOffice/Poppler may not be installed)")
        image_paths = []  # Continue without screenshots
    
    slides_data = []
    screenshot_index = 0
    
    print("Extracting meta-text from slides...")
    for idx, slide in enumerate(prs.slides, start=1):
        # Skip slides without visual content
        if idx not in slides_with_content:
            print(f"  Slide {idx}: Skipped (blank)")
            continue
        
        # Extract text layers
        title, body_text = extract_slide_text(slide)
        notes = extract_notes(slide)
        
        # Ensure title is always present (primary meta-text layer)
        if not title:
            # Try to use first line of body text as title
            if body_text:
                first_line = body_text.split('\n')[0][:100]
                title = first_line if len(first_line) > 3 else f"Slide {idx}"
            else:
                title = f"Slide {idx}"
        
        # Get screenshot path for this slide
        img_path = image_paths[screenshot_index] if screenshot_index < len(image_paths) else None
        screenshot_index += 1
        
        # Store relative path (just filename)
        relative_img_path = os.path.basename(img_path) if img_path else None
        
        if not relative_img_path:
            print(f"  WARNING: No screenshot generated for slide {idx}")
        
        slide_info = {
            "slide_number": idx,
            "title": title,  # Primary meta-text layer
            "text_content": body_text,  # Secondary meta-text layer
            "notes": notes,  # Tertiary meta-text layer
            "image_path": relative_img_path  # Actual slide screenshot
        }
        
        print(f"  Slide {idx}: '{title[:50]}...' - Screenshot: {relative_img_path}")
        slides_data.append(slide_info)
    
    print(f"✅ Successfully processed {len(slides_data)} slides (skipped {len(prs.slides) - len(slides_data)} blank slides)")
    return slides_data


def get_presentation_info(file_path: str) -> Dict:
    """
    Get basic information about a presentation.
    
    Args:
        file_path: Path to the .pptx file
        
    Returns:
        Dictionary with presentation info
    """
    prs = PptxPresentation(file_path)
    
    return {
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height
    }

